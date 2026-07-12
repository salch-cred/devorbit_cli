"""Document and image reading MVP."""
from pathlib import Path

def read_document(workspace: str, rel_path: str, max_chars: int = 20000) -> str:
    root=Path(workspace).resolve(); path=(root/rel_path).resolve()
    if root not in path.parents and path!=root: raise PermissionError('Path escapes workspace')
    if not path.exists(): raise FileNotFoundError(rel_path)
    ext=path.suffix.lower(); text=''
    if ext=='.pdf':
        from pypdf import PdfReader
        text='\n'.join((p.extract_text() or '') for p in PdfReader(str(path)).pages)
    elif ext=='.docx':
        from docx import Document
        text='\n'.join(p.text for p in Document(str(path)).paragraphs)
    elif ext=='.pptx':
        from pptx import Presentation
        text='\n'.join(shape.text for slide in Presentation(str(path)).slides for shape in slide.shapes if hasattr(shape,'text'))
    elif ext in ('.xlsx','.xlsm'):
        from openpyxl import load_workbook
        wb=load_workbook(path,read_only=True,data_only=True); rows=[]
        for ws in wb.worksheets:
            rows.append('['+ws.title+']'); rows.extend('\t'.join('' if v is None else str(v) for v in row) for row in ws.iter_rows(values_only=True))
        text='\n'.join(rows)
    elif ext in ('.png','.jpg','.jpeg','.webp','.bmp','.tiff'):
        import pytesseract
        from PIL import Image
        text=pytesseract.image_to_string(Image.open(path))
    else: text=path.read_text(encoding='utf-8',errors='replace')
    return text[:max_chars]+('\n...[truncated]' if len(text)>max_chars else '')
